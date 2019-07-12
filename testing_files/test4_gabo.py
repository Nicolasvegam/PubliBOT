import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches

df = pd.read_excel('test1.xlsx', sheet_name='Hoja1')
df_2 = pd.read_excel('test1.xlsx', sheet_name='Hoja2')

#print("Column headings:")
head = "\nEl documento leído es\n"
print(head)
print(df)

sex_values = df["Sexo"]
per_values = df["Porcentaje"]

head_s = "\nLos valores para sexo son:\n"
print(head_s)

for element in df["Sexo"]:
    print(element)

head_p = "\nLos valores para porcentaje son:\n"
print(head_p)
for element in df["Porcentaje"]:
    print(element)

# Se crea la ppt
prs = Presentation()

#Se genera la primera slide
slide = prs.slides.add_slide(prs.slide_layouts[0])

chart_data = ChartData()
chart_data.categories = [element for element in df["Sexo"]]
chart_data.add_series('Series 1', (element for element in df["Porcentaje"]))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END



# define chart data ---------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = [element for element in df_2["Edad"]]
chart_data.add_series('Series 2', (element for element in df_2["Porcentaje"]))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

prs.save('test.pptx')
