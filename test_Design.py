import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

class PowerPointer:

    def __init__(self):

        self.presentation = Presentation('Template Reporte Campaña.pptx')

    def test(self):

        #Primera slide
        print(dir(self.presentation.slides[0]))
        slide = self.presentation.slides[0]
#        title = slide.shapes.title
#        title.text = "Título del reporte"
        n = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            else:
                if shape.text_frame.text:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = 'test1'
                    run.font.name = 'Proxima Nova'
                    run.font.size = Pt(18)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255,255,255)
            #print()
            # do things with the text frame
            ...

    def first_slide(self):
        #Se genera la primera slide
        #slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        slide = self.presentation.slides[0]
        title = slide.shapes.title


    def add_table(self, values, title, number):

        title_only_slide_layout = self.presentation.slide_layouts[5]
        #slide = self.presentation.slides.add_slide(title_only_slide_layout)
        slide = self.presentation.slides[number]
        shapes = slide.shapes

        shapes.title.text = title

        left = Inches(1.0)
        top = Inches(0.7)
        width = Inches(5.5)
        height = Inches(0.8)

        cols = len(values[0])
        rows = len(values)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(2.0)

        for i in range(0, cols):
            for j in range(0, rows):
                table.cell(j, i).text = str(values[j][i])

    def add_imagen(self, img_path, number):
        #img_path = 'monty-truth.png'

        blank_slide_layout = self.presentation.slide_layouts[6]
        #slide = self.presentation.slides.add_slide(blank_slide_layout)
        slide = self.presentation.slides[number]
        left = Inches(0.5)
        top = Inches(0.2)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)


    def pie_chart(self, categories, values, title, number):

        #slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide = self.presentation.slides[number]
        shapes = slide.shapes
        shapes.title.text = title

        chart_data = ChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', (element for element in values))
        x, y, cx, cy = Inches(2), Inches(1), Inches(6), Inches(4.5)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END


    def bar_chart(self, categories, values, title, number):

        # define chart data ---------------------
        #slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide = self.presentation.slides[number]
        shapes = slide.shapes
        shapes.title.text = title
        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 2', (element for element in values))

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(1), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    def save(self, file_name):

        self.presentation.save(file_name + '.pptx')

#pr = PowerPointer()
#pr.test()
#pr.save("prueba1")
